import PyPDF2
from docx import Document
from pptx import Presentation
import io

def extract_text(uploaded_file):
    """
    Master function to detect file type and extract text.
    """
    file_name = uploaded_file.name.lower()
    
    try:
        if file_name.endswith('.pdf'):
            return extract_from_pdf(uploaded_file)
        elif file_name.endswith('.docx'):
            return extract_from_docx(uploaded_file)
        elif file_name.endswith('.pptx'):
            return extract_from_pptx(uploaded_file)
        elif file_name.endswith('.txt'):
            return extract_from_txt(uploaded_file)
        else:
            return "Error: Unsupported file format."
    except Exception as e:
        return f"Error processing file: {e}"

def extract_from_pdf(uploaded_file):
    pdf_reader = PyPDF2.PdfReader(uploaded_file)
    text = ""
    for page in pdf_reader.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n"
    return text.strip()

def extract_from_docx(uploaded_file):
    doc = Document(uploaded_file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text.strip()

def extract_from_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def extract_from_txt(uploaded_file):
    # Decode as utf-8, fallback to latin-1
    try:
        return uploaded_file.read().decode('utf-8')
    except UnicodeDecodeError:
        uploaded_file.seek(0)
        return uploaded_file.read().decode('latin-1')
